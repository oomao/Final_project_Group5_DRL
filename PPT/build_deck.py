# -*- coding: utf-8 -*-
"""Hermes-DQN presentation rebuild (white minimal + sticky-note style)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------------- palette ----------------
WHITE=RGBColor(0xFF,0xFF,0xFF)
INK  =RGBColor(0x23,0x27,0x2B)
SLATE=RGBColor(0x46,0x50,0x59)
MUTE =RGBColor(0x8C,0x92,0x98)
HAIR =RGBColor(0xDD,0xE1,0xE5)
GHOST=RGBColor(0xEA,0xED,0xF0)
GREEN=RGBColor(0x2E,0x8B,0x57); GREEN_BG=RGBColor(0xE9,0xF4,0xEC); GREEN_LN=RGBColor(0xBE,0xE2,0xCB)
RED  =RGBColor(0xC4,0x46,0x3D); RED_BG=RGBColor(0xFB,0xEA,0xE8); RED_LN=RGBColor(0xF0,0xC6,0xC1)
AMBER=RGBColor(0xD9,0x9A,0x2B)
STICKY=RGBColor(0xFB,0xF1,0xA9); STICKY_LN=RGBColor(0xE7,0xD4,0x72)
BLUE =RGBColor(0x3F,0x6F,0xA8); BLUE_BG=RGBColor(0xEA,0xF1,0xF8); BLUE_LN=RGBColor(0xC6,0xD9,0xED)
PINK =RGBColor(0xC2,0x5B,0x6B); PINK_BG=RGBColor(0xFB,0xEC,0xEF)
CODEBG=RGBColor(0x25,0x2A,0x31); CODETX=RGBColor(0xE8,0xE8,0xE8)

SERIF="Georgia"; ZH="Microsoft JhengHei"; SANS="Calibri"; MONO="Consolas"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

# ---------------- helpers ----------------
def _font(run,size=14,color=INK,bold=False,italic=False,latin=SANS,ea=ZH):
    f=run.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color
    rPr=run._r.get_or_add_rPr()
    for tag,face in (('a:latin',latin),('a:ea',ea),('a:cs',latin)):
        el=rPr.find(qn(tag))
        if el is None:
            el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set('typeface',face)

def tbx(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP,wrap=True):
    b=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=b.text_frame; tf.word_wrap=wrap
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    tf.vertical_anchor=anchor
    return tf

def para(tf,runs,first=False,align=PP_ALIGN.LEFT,sb=0,sa=0,line=None,level=0):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.level=level
    if sb:p.space_before=Pt(sb)
    if sa:p.space_after=Pt(sa)
    if line is not None:p.line_spacing=line
    for it in runs:
        text=it[0]; opt=it[1] if len(it)>1 else {}
        r=p.add_run(); r.text=text; _font(r,**opt)
    return p

def line_text(slide,l,t,w,h,text,size=14,color=INK,bold=False,italic=False,
              align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,latin=SANS,ea=ZH,wrap=True,line=None):
    tf=tbx(slide,l,t,w,h,anchor,wrap)
    para(tf,[(text,dict(size=size,color=color,bold=bold,italic=italic,latin=latin,ea=ea))],
         first=True,align=align,line=line)
    return tf

def box(slide,l,t,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,radius=None,rotation=0,pad=4):
    sp=slide.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if rotation: sp.rotation=rotation
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    tf=sp.text_frame; tf.word_wrap=True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf,m,Pt(pad))
    return sp

def shape_text(sp,lines,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER):
    tf=sp.text_frame; tf.vertical_anchor=anchor
    for i,(text,opt) in enumerate(lines):
        d=dict(size=12,color=INK,latin=SANS,ea=ZH); d.update(opt)
        ln=d.pop('_line',1.0); sa=d.pop('_sa',0); d.pop('_bullet',None)
        para(tf,[(text,d)],first=(i==0),align=align,line=ln,sa=sa)
    return sp

def connector(slide,x1,y1,x2,y2,color=INK,w=1.5,dash=None):
    c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False
    if dash:
        ln=c.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':dash}); ln.append(d)
    return c

def polyline(slide,pts,color=INK,w=2.0):
    for i in range(len(pts)-1):
        connector(slide,pts[i][0],pts[i][1],pts[i+1][0],pts[i+1][1],color,w)

def arrow(slide,x1,y1,x2,y2,color=MUTE,w=2.25):
    c=connector(slide,x1,y1,x2,y2,color,w)
    ln=c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:headEnd'),{'type':'none'}))
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))
    return c

def dot(slide,cx,cy,r,fill):
    return box(slide,cx-r,cy-r,2*r,2*r,fill=fill,shape=MSO_SHAPE.OVAL)

def chip(slide,l,t,text,fill=INK,fg=WHITE,size=11,w=1.8,h=0.34):
    sp=box(slide,l,t,w,h,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    shape_text(sp,[(text,dict(size=size,color=fg,bold=True))])
    return sp

def marker(slide,l,t,w,h,color=STICKY):
    return box(slide,l,t,w,h,fill=color,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.4)

def sticky(slide,l,t,w,h,lines,rotation=-3,size=12,color=SLATE):
    sp=box(slide,l,t,w,h,fill=STICKY,line=STICKY_LN,lw=0.75,rotation=rotation,pad=8)
    tf=sp.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,(text,opt) in enumerate(lines):
        d=dict(size=size,color=color,latin=SANS,ea=ZH); d.update(opt)
        para(tf,[(text,d)],first=(i==0),align=PP_ALIGN.LEFT,line=1.08)
    return sp

def add_slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
    return s

def header(slide,zh_title,kicker=None,kfill=SLATE,kw=2.0):
    box(slide,0.62,0.56,0.12,0.40,fill=AMBER)
    line_text(slide,0.86,0.5,9.6,0.55,zh_title,size=25,color=INK,bold=True,
              latin=SERIF,ea=ZH,anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        chip(slide,12.7-kw,0.56,kicker,fill=kfill,w=kw,h=0.36,size=11)
    connector(slide,0.62,1.28,12.71,1.28,HAIR,1.0)

def footer(slide,n):
    line_text(slide,0.62,7.07,7,0.3,"Hermes-DQN ・ 國立中興大學 資訊管理學研究所",
              size=8,color=MUTE,latin=SANS,ea=ZH)
    line_text(slide,11.8,7.07,0.9,0.3,str(n),size=9,color=MUTE,align=PP_ALIGN.RIGHT)

def divider(num,en,zh):
    s=add_slide()
    tf=tbx(s,0.55,1.3,3.6,4.6,anchor=MSO_ANCHOR.MIDDLE)
    para(tf,[(str(num),dict(size=300,color=GHOST,bold=True,latin=SERIF,ea=SERIF))],first=True,align=PP_ALIGN.CENTER)
    line_text(s,3.6,2.7,9.0,1.1,en,size=50,color=INK,bold=False,latin=SERIF,ea=ZH,anchor=MSO_ANCHOR.BOTTOM)
    line_text(s,3.66,3.95,9.0,0.8,zh,size=26,color=SLATE,bold=True,latin=ZH,ea=ZH,anchor=MSO_ANCHOR.TOP)
    box(s,3.66,3.78,1.4,0.04,fill=AMBER)
    return s

def style_chart(chart,legend=False,val_axis=False,cat_size=11,catcolor=INK):
    chart.has_title=False; chart.has_legend=legend
    if legend:
        chart.legend.position=XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout=False
        chart.legend.font.size=Pt(11); chart.legend.font.name=SANS
    try:
        va=chart.value_axis; va.has_major_gridlines=False; va.visible=val_axis
        if val_axis:
            va.tick_labels.font.size=Pt(9); va.tick_labels.font.color.rgb=MUTE; va.tick_labels.font.name=SANS
            va.format.line.color.rgb=HAIR
    except Exception: pass
    try:
        ca=chart.category_axis; ca.has_major_gridlines=False
        ca.tick_labels.font.size=Pt(cat_size); ca.tick_labels.font.name=SANS; ca.tick_labels.font.color.rgb=catcolor
        ca.format.line.color.rgb=HAIR
    except Exception: pass

def data_labels(plot,fmt='0.0"%"',size=10,bold=True,pos=XL_LABEL_POSITION.OUTSIDE_END,color=INK):
    plot.has_data_labels=True
    dl=plot.data_labels; dl.number_format=fmt; dl.number_format_is_linked=False
    dl.font.size=Pt(size); dl.font.bold=bold; dl.font.name=SANS; dl.font.color.rgb=color
    try: dl.position=pos
    except Exception: pass

# ============================================================
# SLIDE 1 — title
# ============================================================
s=add_slide()
line_text(s,1.0,1.15,11.3,1.2,"Hermes-DQN",size=62,color=INK,bold=False,
          latin=SERIF,ea=SERIF,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
marker(s,3.05,2.62,7.2,0.62,STICKY)
line_text(s,3.05,2.6,7.2,0.66,"記憶擴增之大型語言模型獎勵設計框架",size=23,color=INK,bold=True,
          latin=ZH,ea=ZH,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# mini architecture pipeline
names=["開源 LLM","Memory","AST Buffer"]
xs=[2.15,4.55,6.95]; bw=2.0; by=4.15; bh=0.95
for x,nm in zip(xs,names):
    b=box(s,x,by,bw,bh,fill=WHITE,line=HAIR,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.12)
    shape_text(b,[(nm,dict(size=13,color=SLATE,bold=True))])
for x in xs:
    arrow(s,x+bw,by+bh/2,x+bw+0.4,by+bh/2,MUTE,2.0)
dqn=box(s,9.35,4.0,2.45,1.25,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.12)
shape_text(dqn,[("DQN Agent",dict(size=16,color=WHITE,bold=True))])
sticky(s,10.35,0.5,2.4,1.0,[("Focus",dict(size=11,bold=True,color=AMBER)),
                            ("Introduction & Motivation",dict(size=12,bold=True))],rotation=-4)
line_text(s,1.0,6.35,8,0.4,"陳盛茂 ・ 林仙安 ・ 辛語柔 ・ 陳冠宇",size=14,color=INK,bold=True,latin=ZH,ea=ZH)
line_text(s,1.0,6.78,8,0.4,"國立中興大學 資訊管理學研究所",size=12,color=MUTE,latin=ZH,ea=ZH)

# ============================================================
# SLIDE 2 — outline
# ============================================================
s=add_slide()
header(s,"Outline 報告大綱")
items=[("1","Introduction","研究背景與動機"),
       ("2","Related Work","相關研究探討"),
       ("3","Proposed Scheme","Hermes-DQN 系統架構"),
       ("4","Simulation","實驗設計與結果分析"),
       ("5","Conclusion","結論與未來展望")]
y=1.75
for num,en,zh in items:
    b=box(s,1.4,y,0.6,0.6,fill=GHOST,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.25)
    shape_text(b,[(num,dict(size=20,color=SLATE,bold=True,latin=SERIF))])
    line_text(s,2.3,y-0.02,5.0,0.4,en,size=22,color=INK,bold=True,latin=SERIF,ea=ZH,anchor=MSO_ANCHOR.TOP)
    line_text(s,2.3,y+0.34,6.5,0.3,zh,size=14,color=MUTE,latin=ZH,ea=ZH)
    y+=1.02
sticky(s,9.7,2.2,2.7,2.3,[("本場聚焦",dict(size=12,bold=True,color=AMBER)),
    ("一個反直覺發現：",dict(size=13,bold=True,color=INK)),
    ("「記憶」對 LLM 獎勵設計",dict(size=12)),
    ("並非永遠有益——",dict(size=12)),
    ("效益取決於獎勵密度。",dict(size=12,bold=True,color=RED))],rotation=3)
footer(s,2)

# ============================================================
# SLIDE 3 — divider 1
# ============================================================
divider(1,"Introduction","研究背景與動機")

# ============================================================
# SLIDE 4 — DRL bottleneck & EUREKA limits
# ============================================================
s=add_slide()
header(s,"DRL 的瓶頸與 EUREKA 的三大侷限",kicker="Motivation",kfill=SLATE,kw=1.9)
# left: two reward problems
lb=box(s,0.7,1.7,2.85,1.6,fill=BLUE_BG,line=BLUE_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
shape_text(lb,[("獎勵設計兩難",dict(size=13,bold=True,color=BLUE,_sa=4)),
               ("稀疏 (Sparse) → 難以學習",dict(size=11.5,color=INK,_sa=2)),
               ("人工塑形 (Dense) → 引入偏差",dict(size=11.5,color=INK))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
# middle: EUREKA
mb=box(s,3.95,1.7,2.7,1.6,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
shape_text(mb,[("EUREKA (ICLR 2024)",dict(size=13,bold=True,color=GREEN,_sa=4)),
               ("以 GPT-4 自動撰寫獎勵函數",dict(size=11.5,color=INK,_sa=2)),
               ("83% 任務超越人類專家",dict(size=12,bold=True,color=INK))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
arrow(s,3.55,2.5,3.95,2.5,MUTE,2.0)
arrow(s,6.65,2.5,7.15,2.5,MUTE,2.0)
# right: three limitations
lims=[("依賴商業 API","GPT-4 費用高、無法離線、不可重現"),
      ("缺乏跨輪記憶","每次迭代從零開始，無法累積經驗"),
      ("忽略緩衝區失效","獎勵變動使 DQN 發生災難性遺忘")]
y=1.7
for t,d in lims:
    rb=box(s,7.15,y,5.45,1.18,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    dot(s,7.55,y+0.36,0.13,RED)
    line_text(s,7.85,y+0.14,4.5,0.4,t,size=14,color=RED,bold=True,latin=ZH,ea=ZH)
    line_text(s,7.85,y+0.58,4.6,0.5,d,size=11,color=SLATE,latin=ZH,ea=ZH,line=1.0)
    y+=1.35
line_text(s,0.7,3.7,6.0,0.6,"→ 三大缺口同時存在，尚無單一框架完整解決。",
          size=13,color=INK,bold=True,latin=ZH,ea=ZH)
sticky(s,0.7,4.55,5.9,1.9,[("研究契機",dict(size=12,bold=True,color=AMBER)),
    ("若能用開源模型 + 記憶 + 緩衝區管理",dict(size=13,color=INK)),
    ("一次補上這三個缺口，",dict(size=13,color=INK)),
    ("自動化獎勵設計才可能真正落地。",dict(size=13,bold=True,color=INK))],rotation=-2)
footer(s,4)

# ============================================================
# SLIDE 5 — research questions & hypothesis
# ============================================================
s=add_slide()
header(s,"核心研究問題與假說",kicker="Research Questions",kfill=SLATE,kw=2.6)
center=box(s,4.15,1.9,5.0,2.2,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(center,[("核心問題",dict(size=12,bold=True,color=STICKY,_sa=6)),
   ("同時解決「開源化、記憶化、",dict(size=15,color=WHITE,bold=True,_line=1.15)),
   ("緩衝區穩定化」之後，",dict(size=15,color=WHITE,bold=True,_line=1.15)),
   ("LLM 設計的獎勵函數",dict(size=15,color=WHITE,bold=True,_line=1.15)),
   ("能否普遍有效？",dict(size=16,color=STICKY,bold=True))],anchor=MSO_ANCHOR.MIDDLE)
q=[("開源化","能否以輕量級開源模型取代 GPT-4？",BLUE,BLUE_BG,BLUE_LN,0.7,2.05),
   ("記憶化","跨迭代累積經驗是否「必定」帶來正向效益？",GREEN,GREEN_BG,GREEN_LN,9.6,2.05),
   ("緩衝穩定","如何在獎勵非穩態下保護 DQN 歷史樣本？",AMBER,RED_BG,STICKY_LN,4.9,4.5)]
for t,d,c,bg,ln,x,y in q:
    b=box(s,x,y,3.0,1.3,fill=bg,line=ln,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    shape_text(b,[(t,dict(size=14,bold=True,color=c,_sa=3)),(d,dict(size=11,color=SLATE,_line=1.05))],
               anchor=MSO_ANCHOR.MIDDLE)
sticky(s,9.4,4.4,3.2,1.55,[("預期 (假說)",dict(size=12,bold=True,color=AMBER)),
    ("記憶機制應能幫 LLM 寫出更好",dict(size=12,color=INK)),
    ("的獎勵，且效益跨任務一致。",dict(size=12,color=INK)),
    ("→ 本研究將檢驗它。",dict(size=12,bold=True,color=RED))],rotation=2)
footer(s,5)

# ============================================================
# SLIDE 6 — divider 2
# ============================================================
divider(2,"Related Work","相關研究探討")

# ============================================================
# SLIDE 7 — Pillar 1 table
# ============================================================
s=add_slide()
header(s,"Pillar 1｜LLM 撰寫獎勵函數之演進與侷限",kicker="Related Work",kfill=SLATE,kw=2.0)
cols=[("","",2.7),("核心模型 / API","",2.9),("跨代記憶 (Memory)","",2.9),("緩衝區遺忘處理","",3.1)]
x0=0.7; y0=1.65; rowh=0.78
xacc=x0; colx=[]
for _,_,w in cols:
    colx.append((xacc,w)); xacc+=w
# header row
for (title_,_,w),(cx,_) in zip(cols,colx):
    if title_:
        line_text(s,cx+0.1,y0+0.1,w-0.2,0.5,title_,size=12.5,color=SLATE,bold=True,latin=ZH,ea=ZH,anchor=MSO_ANCHOR.MIDDLE)
rows=[("EUREKA","(Ma et al., ICLR 2024)","依賴 GPT-4","無記憶","無處理"),
      ("CARD","(Sun et al., 2024)","依賴 LLM 評論員","無原生環境記憶","無處理"),
      ("LEARN-Opt","(Cardenoso & Caarls, 2025)","開源小模型","無記憶","無處理")]
y=y0+0.7
for nm,cite,c1,c2,c3 in rows:
    connector(s,x0,y,x0+11.6,y,HAIR,1.0)
    line_text(s,colx[0][0]+0.1,y+0.08,colx[0][1]-0.2,0.35,nm,size=13,color=INK,bold=True,latin=SERIF,ea=ZH)
    line_text(s,colx[0][0]+0.1,y+0.42,colx[0][1]-0.2,0.3,cite,size=9.5,color=MUTE,latin=SANS,ea=ZH)
    for val,(cx,w) in zip((c1,c2,c3),colx[1:]):
        line_text(s,cx+0.1,y+0.1,w-0.2,0.55,val,size=11.5,color=SLATE,latin=ZH,ea=ZH,anchor=MSO_ANCHOR.MIDDLE)
    y+=rowh
connector(s,x0,y,x0+11.6,y,HAIR,1.0)
gap=box(s,0.7,4.85,11.9,1.55,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
line_text(s,1.0,5.0,4.0,0.4,"Research Gap｜當前三大侷限",size=14,color=GREEN,bold=True,latin=ZH,ea=ZH)
g=["過度依賴昂貴且無法離線的商業 API","缺乏跨迭代的經驗累積 (Memory)","未處理獎勵變更導致的歷史樣本失效"]
gx=1.0
for i,t in enumerate(g):
    line_text(s,gx,5.5,3.7,0.8,"%d. %s"%(i+1,t),size=12,color=INK,bold=True,latin=ZH,ea=ZH,line=1.1)
    gx+=3.85
footer(s,7)

# ============================================================
# SLIDE 8 — Pillar 2 memory hierarchy
# ============================================================
s=add_slide()
header(s,"Pillar 2｜記憶擴增之 LLM 智能體",kicker="Related Work",kfill=SLATE,kw=2.0)
line_text(s,0.7,1.5,7.6,0.7,"「記憶架構是近期智能體進步的最大推手：OSWorld 準確率 12% → 66.3%」",
          size=12.5,color=SLATE,italic=True,latin=ZH,ea=ZH,line=1.15)
line_text(s,0.7,2.1,7.6,0.3,"— Stanford HAI 2026 AI Index",size=10.5,color=MUTE,latin=SANS,ea=ZH)
layers=[("Procedural Memory","SKILL.md — 固化操作規則"),
        ("Semantic Memory","USER / MEMORY.md — 成功/失敗範例知識庫"),
        ("Episodic Memory","SQLite FTS5 — 完整歷史與 Fitness 指標"),
        ("Working Memory","Prompt Context — 當前檢索的 Top-K 先驗")]
y=2.65
for en,zh in layers:
    b=box(s,0.7,y,7.4,0.82,fill=BLUE_BG,line=BLUE_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    line_text(s,0.95,y+0.12,3.0,0.3,en,size=13,color=BLUE,bold=True,latin=SERIF,ea=ZH)
    line_text(s,0.95,y+0.45,7.0,0.3,zh,size=11,color=SLATE,latin=ZH,ea=ZH)
    y+=0.96
arrow(s,8.4,3.4,9.2,3.4,MUTE,2.5)
rg=box(s,9.0,2.3,3.6,2.9,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(rg,[("Research Gap",dict(size=14,bold=True,color=RED,_sa=4)),
   ("效益的未定之天",dict(size=15,bold=True,color=INK,_sa=8)),
   ("記憶對 Agent 有益，但對「DQN 獎勵設計」是否有絕對助益？",dict(size=11.5,color=SLATE,_line=1.2,_sa=6)),
   ("尤其密集 vs 稀疏環境的差異，至今缺乏嚴謹的 DRL 多任務消融實證。",dict(size=11.5,color=SLATE,_line=1.2))],
   anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,8)

# ============================================================
# SLIDE 9 — Pillar 3 Bellman drift
# ============================================================
s=add_slide()
header(s,"Pillar 3｜非穩態獎勵與底層穩定性 (Bellman Drift)",kicker="Related Work",kfill=SLATE,kw=2.0)
flow=["LLM 變更獎勵函數","目標分佈飄移 (Ng et al., 1999)","Bellman 算子漂移 / CHAIN 效應","歷史樣本失效 (災難性遺忘)"]
y=1.65
for i,t in enumerate(flow):
    fill=RED_BG if i==len(flow)-1 else GHOST
    ln=RED_LN if i==len(flow)-1 else HAIR
    tc=RED if i==len(flow)-1 else INK
    b=box(s,0.7,y,5.2,0.85,fill=fill,line=ln,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    shape_text(b,[(t,dict(size=13,color=tc,bold=(i==len(flow)-1)))],align=PP_ALIGN.LEFT)
    if i<len(flow)-1: arrow(s,3.3,y+0.85,3.3,y+1.0,MUTE,2.0)
    y+=1.0
sol=box(s,6.4,1.65,6.2,1.55,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(sol,[("過去解法 (Value-side)",dict(size=13,bold=True,color=GREEN,_sa=4)),
   ("GB-DQN (Lee & Lee, 2025)：以梯度增強",dict(size=12,color=INK,_line=1.15)),
   ("處理「價值函數面」的遺忘。",dict(size=12,color=INK))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
gap=box(s,6.4,3.45,6.2,2.55,fill=STICKY,line=STICKY_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
shape_text(gap,[("Research Gap",dict(size=14,bold=True,color=AMBER,_sa=4)),
   ("缺乏直接針對「重播緩衝區 (Buffer-side)」的處理。",dict(size=13,bold=True,color=INK,_line=1.2,_sa=6)),
   ("→ 需建立一套依「程式碼差異 (AST Diff)」決定樣本保留策略的底層機制，",dict(size=12,color=SLATE,_line=1.2)),
   ("與 value-side 工作正交互補。",dict(size=12,color=SLATE))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,9)

# ============================================================
# SLIDE 10 — synthesis puzzle mapping
# ============================================================
s=add_slide()
header(s,"Synthesis｜Hermes-DQN 的研究定位",kicker="Positioning",kfill=SLATE,kw=1.9)
line_text(s,0.7,1.4,5.0,0.4,"文獻缺口 (Limitations)",size=13,color=MUTE,bold=True,latin=ZH,ea=ZH)
line_text(s,7.0,1.4,5.6,0.4,"Hermes-DQN 架構創新",size=13,color=GREEN,bold=True,latin=ZH,ea=ZH)
pairs=[("依賴閉源 API、無記憶累積","開源 LLM 替代 + 導入四層記憶閉環架構"),
       ("DRL 任務中記憶機制缺乏實證","首度發現記憶在密集環境的負向效應 (p=0.0317)，並以 DQN 變體驗證模型無關性"),
       ("緩衝區樣本遺忘問題未解","獨創 AST 感知重播緩衝管理，依結構差異執行 KEEP / PARTIAL / DECAY / CLEAR")]
y=1.95
for lim,inn in pairs:
    lb=box(s,0.7,y,5.5,1.15,fill=GHOST,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(lb,[(lim,dict(size=12.5,color=INK,bold=True,_line=1.15))],align=PP_ALIGN.LEFT)
    arrow(s,6.3,y+0.57,6.85,y+0.57,GREEN,2.5)
    rb=box(s,6.95,y,5.65,1.15,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(rb,[(inn,dict(size=11.5,color=INK,_line=1.2))],align=PP_ALIGN.LEFT)
    y+=1.32
b=box(s,0.7,5.95,11.9,0.7,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
shape_text(b,[("Hermes-DQN 整合三大領域，建立開源自動化獎勵設計的研究基準。",dict(size=14,bold=True,color=WHITE))])
footer(s,10)

# ============================================================
# SLIDE 11 — divider 3
# ============================================================
divider(3,"Proposed Scheme","Hermes-DQN 系統架構")

# ============================================================
# SLIDE 12 — architecture big picture
# ============================================================
s=add_slide()
header(s,"系統架構資料流 (The Big Picture)",kicker="Architecture",kfill=SLATE,kw=1.9)
sub=[("子系統一：跨迭代學習能力",BLUE),("子系統二：開源且可重現的獎勵引擎",GREEN),("子系統三：抑制漂移與災難性遺忘",AMBER)]
sx=0.7
for t,c in sub:
    dot(s,sx+0.1,1.65,0.08,c); line_text(s,sx+0.28,1.5,3.9,0.35,t,size=11,color=SLATE,bold=True,latin=ZH,ea=ZH); sx+=4.05
blocks=[("Long-term\nMemory","SQLite FTS5",BLUE_BG,BLUE_LN,BLUE),
        ("開源 LLM","reward author",GREEN_BG,GREEN_LN,GREEN),
        ("AST diff\n+ Buffer","KEEP/PARTIAL/\nDECAY/CLEAR",STICKY,STICKY_LN,AMBER),
        ("DQN Agent","the learner",PINK_BG,RED_LN,PINK),
        ("Env-native\nEval","100 unseen seeds",GHOST,HAIR,SLATE)]
bw=2.05; gap=0.45; bx=0.55; by=2.7; bh=1.9
labels=["(1) priors","(2) reward.py","(3) buffer","(4) model"]
xpos=[]
for i,(t,sub2,bg,ln,c) in enumerate(blocks):
    x=bx+i*(bw+gap); xpos.append(x)
    b=box(s,x,by,bw,bh,fill=bg,line=ln,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    if i==3:
        shape_text(b,[(t,dict(size=14,bold=True,color=c,_line=1.0,_sa=4)),
                      ("Q-net 64×64 MLP",dict(size=9,color=SLATE,_line=1.05)),
                      ("Target net (lagged)",dict(size=9,color=SLATE,_line=1.05)),
                      ("Replay 100K · ε1.0→0.01",dict(size=9,color=SLATE,_line=1.05))],anchor=MSO_ANCHOR.MIDDLE)
    else:
        shape_text(b,[(t,dict(size=14,bold=True,color=c,_line=1.0,_sa=4)),(sub2,dict(size=9.5,color=SLATE,_line=1.05))],anchor=MSO_ANCHOR.MIDDLE)
for i in range(4):
    ax=xpos[i]+bw; arrow(s,ax,by+bh/2,ax+gap,by+bh/2,MUTE,2.0)
    line_text(s,ax-0.15,by-0.32,bw,0.3,labels[i],size=9,color=MUTE,align=PP_ALIGN.CENTER,latin=SANS,ea=ZH)
# return arrow
ry=by+bh+0.55
connector(s,xpos[4]+bw/2,by+bh,xpos[4]+bw/2,ry,MUTE,2.0)
arrow(s,xpos[4]+bw/2,ry,xpos[0]+bw/2,ry,MUTE,2.0)
connector(s,xpos[0]+bw/2,ry,xpos[0]+bw/2,by+bh,MUTE,2.0)
line_text(s,2.9,ry-0.02,7.5,0.3,"(5) fitness writeback — 閉環迴路（每組環境/條件/種子重複 5 次）",size=10,color=MUTE,align=PP_ALIGN.CENTER,latin=SANS,ea=ZH)
footer(s,12)

# ============================================================
# SLIDE 13 — 7 step loop
# ============================================================
s=add_slide()
header(s,"7 步驟迭代閉環：系統運轉引擎",kicker="Closed Loop",kfill=SLATE,kw=1.9)
steps=[("1 檢索","memory.top_k_by_fitness()\n取前 5 筆最佳歷史"),
       ("2 生成","LLM 接任務規格 + 先驗\n生成 reward_src"),
       ("3 差異","解析新舊規格之\nAST 差異程度"),
       ("4 緩衝","依差異決定保留\n或清空 prev_buffer"),
       ("5 訓練","DQN 執行 N 輪\nepisodes 訓練"),
       ("6 評估","100 個未見種子\n原生獎勵公平評估"),
       ("7 寫回","(原始碼, 適應度)\n寫回記憶體")]
positions=[(0.7,1.7),(3.65,1.7),(6.6,1.7),(9.55,1.7),(9.55,4.05),(6.6,4.05),(3.65,4.05)]
bw=2.7; bh=1.55
for i,((t,d),(x,y)) in enumerate(zip(steps,positions)):
    hl = (i==5)
    b=box(s,x,y,bw,bh,fill=(GREEN_BG if hl else WHITE),line=(GREEN_LN if hl else HAIR),lw=1.25,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(b,[(t,dict(size=14,bold=True,color=(GREEN if hl else INK),_sa=4)),
                  (d,dict(size=10.5,color=SLATE,_line=1.1))],anchor=MSO_ANCHOR.MIDDLE)
# arrows along the snake
arrow(s,3.4,2.47,3.65,2.47); arrow(s,6.35,2.47,6.6,2.47); arrow(s,9.3,2.47,9.55,2.47)
arrow(s,10.9,3.25,10.9,4.05)  # down 4->5
arrow(s,9.55,4.82,9.3,4.82); arrow(s,6.6,4.82,6.35,4.82); arrow(s,3.65,4.82,3.4,4.82)
sticky(s,0.7,4.05,2.7,1.55,[("關鍵設計",dict(size=11,bold=True,color=AMBER)),
   ("第 6 步固定用 100 個",dict(size=11,color=INK)),
   ("未見種子 + 原生獎勵，",dict(size=11,color=INK)),
   ("是跨條件唯一公平基準。",dict(size=11,bold=True,color=INK))],rotation=-2)
footer(s,13)

# ============================================================
# SLIDE 14 — 4-layer memory
# ============================================================
s=add_slide()
header(s,"放大檢視 I：Hermes 四層記憶架構",kicker="Zoom-in I",kfill=BLUE,kw=1.8)
layers=[("Working Memory","作用中記憶","Prompt Context（本輪檢索出的 top-K 先驗）"),
        ("Episodic Memory","情節記憶","SQLite FTS5（Fitness 指標、獎勵原始碼、差異分類）"),
        ("Semantic Memory","語義記憶","MEMORY.md（長期成功 / 失敗教訓）"),
        ("Procedural Memory","程序記憶","SKILL.md（固化操作規則，如要求純函數）")]
y=1.7
for en,zh,desc in layers:
    b=box(s,0.7,y,4.5,1.0,fill=BLUE_BG,line=BLUE_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(b,[(en,dict(size=14,bold=True,color=BLUE,_sa=2)),(zh,dict(size=11,color=SLATE))],
               anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    line_text(s,5.45,y+0.05,4.3,0.9,desc,size=12,color=INK,latin=ZH,ea=ZH,anchor=MSO_ANCHOR.MIDDLE,line=1.15)
    y+=1.18
# retrieval funnel
fb=box(s,10.0,1.9,2.6,2.4,fill=GHOST,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(fb,[("檢索漏斗",dict(size=13,bold=True,color=INK,_sa=6)),
   ("Retrieval Funnel",dict(size=10,color=MUTE,_sa=8)),
   ("查詢鍵 =",dict(size=11.5,color=SLATE,_line=1.2)),
   ("環境名稱 + 上一輪",dict(size=11.5,color=SLATE,_line=1.2)),
   ("Fitness 指標",dict(size=11.5,color=SLATE,_line=1.2)),
   ("→ 取 Top-K",dict(size=12,bold=True,color=BLUE))],anchor=MSO_ANCHOR.MIDDLE)
footer(s,14)

# ============================================================
# SLIDE 15 — Gemma reward generator
# ============================================================
s=add_slide()
header(s,"放大檢視 II：開源獎勵生成器",kicker="Zoom-in II",kfill=GREEN,kw=1.8)
ins=["Task Spec（環境描述、觀察空間）","歷史先驗（來自四層記憶）","Fitness 回饋"]
line_text(s,0.7,1.55,3.6,0.3,"輸入 (Input)",size=12,color=MUTE,bold=True,latin=SANS,ea=ZH)
y=1.95
for t in ins:
    b=box(s,0.7,y,3.6,0.7,fill=CODEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    shape_text(b,[(t,dict(size=11.5,color=CODETX))],align=PP_ALIGN.LEFT)
    y+=0.85
arrow(s,4.4,2.9,5.0,2.9,MUTE,2.5)
eng=box(s,5.05,1.95,3.4,1.9,fill=GREEN_BG,line=GREEN_LN,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
shape_text(eng,[("引擎：開源 LLM",dict(size=16,bold=True,color=GREEN,_sa=6)),
   ("Stateless 無狀態：",dict(size=11.5,color=INK,_line=1.15)),
   ("所有跨迭代訊息",dict(size=11.5,color=INK,_line=1.15)),
   ("皆由記憶區塊載入",dict(size=11.5,color=INK,_line=1.15))],anchor=MSO_ANCHOR.MIDDLE)
arrow(s,8.55,2.9,9.15,2.9,MUTE,2.5)
line_text(s,9.2,1.55,3.4,0.3,"輸出 (Output)",size=12,color=MUTE,bold=True,latin=SANS,ea=ZH)
outb=box(s,9.2,1.95,3.4,1.9,fill=CODEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(outb,[("def reward_fn(",dict(size=11.5,color=CODETX,latin=MONO,_line=1.25)),
   ("  obs, action,",dict(size=11.5,color=CODETX,latin=MONO,_line=1.25)),
   ("  next_obs, done",dict(size=11.5,color=CODETX,latin=MONO,_line=1.25)),
   (") -> float",dict(size=11.5,color=STICKY,latin=MONO,_line=1.25))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
sb=box(s,0.7,4.55,11.9,1.6,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(sb,[("安全沙箱 (Security Sandbox)",dict(size=14,bold=True,color=RED,_sa=6)),
   ("語法 AST 解析 + 沙箱執行驗證（L2 subprocess isolation）：僅允許決定性 Python，禁止網路與檔案 I/O。",
    dict(size=12.5,color=INK,_line=1.2))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,15)

# ============================================================
# SLIDE 16 — AST-aware buffer spectrum
# ============================================================
s=add_slide()
header(s,"放大檢視 III：AST 感知 Replay Buffer 策略",kicker="Zoom-in III",kfill=AMBER,kw=1.8)
line_text(s,0.7,1.5,11.9,0.4,"相似度決策光譜 (Similarity Decision Spectrum)",size=13,color=SLATE,bold=True,
          latin=ZH,ea=ZH,align=PP_ALIGN.CENTER)
# gradient bar (green -> red) via 4 segments
segx=[1.1,4.0,6.9,9.8]; segw=2.85; bary=2.5; barh=0.5
segc=[GREEN,RGBColor(0x9C,0xC4,0x6A),AMBER,RED]
for x,c in zip(segx,segc):
    box(s,x,bary,segw,barh,fill=c,shape=MSO_SHAPE.RECTANGLE)
line_text(s,0.55,bary+0.6,1.0,0.3,"1.00",size=11,color=GREEN,bold=True,align=PP_ALIGN.CENTER,latin=SANS)
line_text(s,11.85,bary+0.6,1.0,0.3,"0.00",size=11,color=RED,bold=True,align=PP_ALIGN.CENTER,latin=SANS)
zones=[("IDENTICAL","KEEP","完整保留",GREEN),
       ("SIGNATURE_ONLY","PARTIAL_KEEP","保留高 Q 樣本",RGBColor(0x6F,0x9E,0x3F)),
       ("STRUCTURAL_DIFF","DECAY","權重 ×0.5 衰減",AMBER),
       ("TOTAL_REWRITE","CLEAR","清空重來",RED)]
for (cat,act,desc,c),x in zip(zones,segx):
    arrow(s,x+segw/2,bary+barh,x+segw/2,bary+barh+0.25,c,2.0)
    line_text(s,x-0.1,bary+0.78,segw+0.2,0.3,cat,size=10,color=MUTE,align=PP_ALIGN.CENTER,latin=MONO,ea=ZH)
    b=box(s,x+0.2,bary+1.12,segw-0.4,0.95,fill=WHITE,line=c,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.12)
    shape_text(b,[(act,dict(size=14,bold=True,color=c,latin=MONO,_sa=2)),(desc,dict(size=10.5,color=SLATE))],anchor=MSO_ANCHOR.MIDDLE)
line_text(s,0.7,5.5,11.9,0.5,"差異越大，清除越狠",size=18,color=INK,bold=True,latin=ZH,ea=ZH,align=PP_ALIGN.CENTER)
line_text(s,0.7,6.05,11.9,0.4,"Linear degradation of prior experience based on structural drift",
          size=12,color=MUTE,italic=True,align=PP_ALIGN.CENTER,latin=SANS)
footer(s,16)

# ============================================================
# SLIDE 17 — AST diff computation (NEW)
# ============================================================
s=add_slide()
header(s,"AST 差異如何計算？（補充）",kicker="How it works",kfill=AMBER,kw=2.0)
codeb=box(s,0.7,1.6,5.5,2.4,fill=CODEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
shape_text(codeb,[("# 兩段原始碼字串的純函數",dict(size=11,color=MUTE,latin=MONO,_line=1.3)),
   ("old = ast.parse(prev_src)",dict(size=11.5,color=CODETX,latin=MONO,_line=1.3)),
   ("new = ast.parse(curr_src)",dict(size=11.5,color=CODETX,latin=MONO,_line=1.3)),
   ("kind = classify(old, new)",dict(size=11.5,color=STICKY,latin=MONO,_line=1.3)),
   ("policy = DECIDE[kind]",dict(size=11.5,color=STICKY,latin=MONO,_line=1.3)),
   ("apply(prev_buffer, policy)",dict(size=11.5,color=CODETX,latin=MONO,_line=1.3))],
   anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
line_text(s,0.7,4.15,5.6,1.6,
   "以 Python 內建 ast 模組分別解析新舊獎勵原始碼，比較語法樹結構；分類器與策略查表皆為「兩個字串的決定性純函數」，可重現、無副作用。",
   size=12.5,color=SLATE,latin=ZH,ea=ZH,line=1.3)
line_text(s,6.6,1.55,6.0,0.4,"四類結構差異 → 對應策略",size=13,color=INK,bold=True,latin=ZH,ea=ZH)
cls=[("IDENTICAL","語法樹完全相同","KEEP",GREEN),
     ("SIGNATURE_ONLY","簽章不變、函式體變動","PARTIAL_KEEP",RGBColor(0x6F,0x9E,0x3F)),
     ("STRUCTURAL_DIFF","控制流變動、多數運算元延續","DECAY",AMBER),
     ("TOTAL_REWRITE","結構幾無重疊","CLEAR",RED)]
y=2.05
for cat,desc,act,c in cls:
    b=box(s,6.6,y,6.0,0.92,fill=WHITE,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    box(s,6.6,y,0.12,0.92,fill=c)
    line_text(s,6.9,y+0.1,3.5,0.3,cat,size=12,color=c,bold=True,latin=MONO,ea=ZH)
    line_text(s,6.9,y+0.45,4.3,0.4,desc,size=10.5,color=SLATE,latin=ZH,ea=ZH)
    line_text(s,11.2,y+0.12,1.3,0.7,"→ "+act,size=11.5,color=c,bold=True,latin=MONO,ea=ZH,anchor=MSO_ANCHOR.MIDDLE)
    y+=1.0
footer(s,17)

# ============================================================
# SLIDE 18 — DQN variants (model agnostic)
# ============================================================
s=add_slide()
header(s,"智能體變體：展現系統的模型無關性",kicker="Model-Agnostic",kfill=SLATE,kw=2.3)
line_text(s,0.7,1.4,11.9,0.4,"Hermes 框架與底層 DQN 架構完全解耦，可透過 --dqn-variant 無縫切換。",
          size=13,color=SLATE,latin=ZH,ea=ZH)
cells=[("1  Vanilla","use_double=False\ndueling=False","標準 DQN",GHOST,HAIR,INK),
       ("2  Double DQN","use_double=True\ndueling=False","分離線上/目標網路，消除最大化偏誤",BLUE_BG,BLUE_LN,BLUE),
       ("3  Dueling DQN","use_double=False\ndueling=True","Q 分解為 V(s) 與優勢 A(s,a)",GREEN_BG,GREEN_LN,GREEN),
       ("4  Double + Dueling","use_double=True\ndueling=True","兩者結合",STICKY,STICKY_LN,AMBER)]
pos=[(0.7,1.95),(6.65,1.95),(0.7,4.05),(6.65,4.05)]
cw=5.95; ch=1.9
for (t,code,d,bg,ln,c),(x,y) in zip(cells,pos):
    b=box(s,x,y,cw,ch,fill=bg,line=ln,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
    line_text(s,x+0.3,y+0.18,cw-0.6,0.4,t,size=16,color=c,bold=True,latin=SERIF,ea=ZH)
    cb=box(s,x+0.3,y+0.68,2.6,0.95,fill=CODEBG,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(cb,[(code.split("\n")[0],dict(size=10.5,color=CODETX,latin=MONO,_line=1.2)),
                   (code.split("\n")[1],dict(size=10.5,color=CODETX,latin=MONO,_line=1.2))],align=PP_ALIGN.LEFT)
    line_text(s,x+3.1,y+0.7,cw-3.3,0.95,d,size=11.5,color=SLATE,latin=ZH,ea=ZH,anchor=MSO_ANCHOR.MIDDLE,line=1.2)
footer(s,18)

# ============================================================
# SLIDE 19 — ablation matrix
# ============================================================
s=add_slide()
header(s,"嚴謹驗證：六種消融實驗條件 (Ablation Matrix)",kicker="Ablation",kfill=SLATE,kw=1.7)
heads=["條件 (Condition)","獎勵來源","記憶 (FTS5)","AST 緩衝","迭代"]
colx=[0.7,3.9,6.3,8.4,10.5]; colw=[3.2,2.4,2.1,2.1,2.1]
y0=1.7
for h,x,w in zip(heads,colx,colw):
    line_text(s,x+0.05,y0,w-0.1,0.4,h,size=12,color=SLATE,bold=True,latin=ZH,ea=ZH)
rows=[("B0-env-native","原生獎勵","—","—","1","Baseline 基準線",False),
      ("B1-handcrafted","人工塑形","—","—","1","人類專家基準",False),
      ("B2-gemma-oneshot","LLM 單次","—","—","1","LLM 基準",False),
      ("B3-no-memory","LLM","—","✓","5","驗證記憶的獨立影響",False),
      ("B3-no-AST","LLM","✓","—","5","驗證緩衝管理的獨立影響",False),
      ("B3-hermes-full","LLM","✓","✓","5","完整系統",True)]
y=y0+0.45
for nm,src,mem,ast,it,note,hl in rows:
    if hl:
        box(s,0.65,y-0.05,11.95,0.66,fill=STICKY,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    else:
        connector(s,0.7,y+0.58,12.6,y+0.58,HAIR,0.75)
    line_text(s,colx[0]+0.05,y+0.04,colw[0],0.4,nm,size=12.5,color=INK,bold=hl,latin=MONO,ea=ZH)
    for val,x,w in zip((src,mem,ast,it),colx[1:],colw[1:]):
        col = GREEN if val=="✓" else (MUTE if val=="—" else INK)
        line_text(s,x+0.05,y+0.04,w-0.1,0.4,val,size=12,color=col,bold=(val=="✓"),latin=ZH,ea=ZH)
    line_text(s,colx[1]+0.05,y+0.33,8,0.25,note,size=9.5,color=MUTE,latin=ZH,ea=ZH)
    y+=0.64
sticky(s,8.4,6.08,4.2,0.82,[("總規模：4 環境 × 6 條件 × 5 種子",dict(size=11,bold=True,color=INK)),
   ("= 120 次完整訓練",dict(size=12.5,color=RED,bold=True))],rotation=-2)
footer(s,19)

# ============================================================
# SLIDE 20 — divider 4
# ============================================================
divider(4,"Simulation","實驗設計與結果分析")

# ============================================================
# SLIDE 21 — experimental setup (NEW)
# ============================================================
s=add_slide()
header(s,"實驗設定 (Setup)",kicker="Setup",kfill=SLATE,kw=1.5)
line_text(s,0.7,1.45,6.0,0.35,"四個古典控制環境（涵蓋獎勵密度兩端）",size=13,color=INK,bold=True,latin=ZH,ea=ZH)
envh=["環境","密度","觀測","動作","解題門檻"]
ex=[0.7,2.65,3.95,4.95,5.95]; ew=[1.95,1.3,1.0,1.0,1.4]
for h,x,w in zip(envh,ex,ew):
    line_text(s,x+0.03,1.85,w,0.3,h,size=10.5,color=MUTE,bold=True,latin=ZH,ea=ZH)
envs=[("LunarLander-v3","密集",RED,"8","4","mean ≥ 200"),
      ("CartPole-v1","稀疏",GREEN,"4","2","mean ≥ 475"),
      ("MountainCar-v0","稀疏",GREEN,"2","3","mean ≥ −110"),
      ("Acrobot-v1","稀疏",GREEN,"6","3","mean ≥ −100")]
y=2.2
for nm,dens,dc,obs,act,thr in envs:
    connector(s,0.7,y+0.5,7.3,y+0.5,HAIR,0.75)
    line_text(s,ex[0]+0.03,y+0.08,ew[0],0.4,nm,size=11.5,color=INK,bold=True,latin=SANS,ea=ZH)
    cl=chip(s,ex[1]+0.03,y+0.08,dens,fill=dc,w=0.9,h=0.32,size=10)
    line_text(s,ex[2]+0.03,y+0.1,ew[2],0.35,obs,size=11.5,color=SLATE)
    line_text(s,ex[3]+0.03,y+0.1,ew[3],0.35,act,size=11.5,color=SLATE)
    line_text(s,ex[4]+0.03,y+0.1,ew[4],0.35,thr,size=11.5,color=SLATE,latin=SANS,ea=ZH)
    y+=0.62
# right config card
cfg=box(s,7.7,1.45,4.9,4.45,fill=GHOST,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
rows=[("硬體","RTX 4090 ×1 · Win11 · CUDA 12.1"),
      ("環境","Python 3.11 · PyTorch 2.5.1"),
      ("DQN","64×64 MLP · lr=5e-4 · γ=0.99"),
      ("","batch=64 · replay 100K"),
      ("","ε 於 50K 步線性衰減 → 0.01"),
      ("","目標網路每 1000 步更新"),
      ("訓練種子","42, 43, 44, 45, 46"),
      ("評估種子","10000–10099（與訓練互斥）"),
      ("總訓練","120 次 = 4×6×5")]
yy=1.7
line_text(s,7.95,yy,4.4,0.35,"訓練配置",size=13,color=INK,bold=True,latin=ZH,ea=ZH); yy+=0.5
for k,v in rows:
    if k: line_text(s,7.95,yy,1.2,0.35,k,size=11,color=BLUE,bold=True,latin=ZH,ea=ZH)
    line_text(s,9.15,yy,3.35,0.35,v,size=11,color=SLATE,latin=SANS,ea=ZH)
    yy+=0.42
footer(s,21)

# ============================================================
# SLIDE 22 — metric + statistical methodology (NEW)
# ============================================================
s=add_slide()
header(s,"評估指標與統計方法",kicker="Methodology",kfill=SLATE,kw=2.0)
mb=box(s,0.7,1.55,11.9,1.25,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(mb,[("主要指標：env_native_mean",dict(size=15,bold=True,color=STICKY,_sa=4)),
   ("以「環境原生（未經塑形）獎勵」在 100 個未見種子 (10000–10099) 上的平均回報——這是跨條件唯一公平的衡量基準。",
    dict(size=12.5,color=WHITE,_line=1.2))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
cards=[("檢定方法","雙尾 Mann-Whitney U\n(α = 0.05，非參數、適合 n=5)",BLUE_BG,BLUE_LN,BLUE),
       ("信賴區間","Bootstrap 5000 次再抽樣\n95% 信賴水準",GREEN_BG,GREEN_LN,GREEN),
       ("「勝出」三條件","p<0.05、|Δ|/|base|≥10%、\n信賴區間不重疊",STICKY,STICKY_LN,AMBER)]
x=0.7
for t,d,bg,ln,c in cards:
    b=box(s,x,3.05,3.83,1.5,fill=bg,line=ln,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    shape_text(b,[(t,dict(size=13,bold=True,color=c,_sa=4)),(d,dict(size=11,color=INK,_line=1.2))],
               anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
    x+=4.03
warn=box(s,0.7,4.8,11.9,1.35,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(warn,[("統計力限制（誠實揭露）",dict(size=13,bold=True,color=RED,_sa=4)),
   ("n=5 僅對「大效應」(Cohen's d ≥ 1) 具穩定偵測力；中等效應在此樣本下難以定論。種子全數保留、不剔除崩潰種子（如 Acrobot 之 B0），故其變異與信賴區間被膨脹。",
    dict(size=12,color=INK,_line=1.25))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,22)

# ============================================================
# SLIDE 23 — Part 1 headline (diverging bar) [orig S20]
# ============================================================
s=add_slide()
header(s,"Part 1 結果：稀疏「大勝」、密集「反轉」",kicker="Part 1 · Results",kfill=SLATE,kw=2.3)
cd=CategoryChartData()
cd.categories=["LunarLander","CartPole","MountainCar","Acrobot"]
cd.add_series("Δ% vs B0",(-11.4,116.1,31.5,57.5))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(0.7),Inches(1.55),Inches(7.4),Inches(4.7),cd)
chart=gf.chart; style_chart(chart,legend=False,val_axis=False,cat_size=12)
data_labels(chart.plots[0],fmt='+0.0"%";-0.0"%"',size=12)
ser=chart.series[0]
for i,c in enumerate((RED,GREEN,GREEN,GREEN)):
    ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb=c
line_text(s,0.7,6.25,7.4,0.3,"B3-hermes-full vs B0-env-native（p：1.00 / 0.0317 / 0.0112 / 0.0952）",
          size=10,color=MUTE,align=PP_ALIGN.CENTER,latin=SANS,ea=ZH)
gbox=box(s,8.35,1.7,4.3,1.95,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(gbox,[("稀疏環境：顯著正向",dict(size=13,bold=True,color=GREEN,_sa=4)),
   ("CartPole +116% (p=0.0317)",dict(size=12,color=INK,_line=1.25)),
   ("MountainCar +31.5% (p=0.0112)",dict(size=12,color=INK,_line=1.25)),
   ("→ 原生訊號微弱時，LLM 解鎖學習瓶頸。",dict(size=11,color=SLATE,_line=1.2))],
   anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
rbox=box(s,8.35,3.85,4.3,2.3,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(rbox,[("密集環境：記憶帶來傷害",dict(size=13,bold=True,color=RED,_sa=4)),
   ("LunarLander：相對「無記憶版」",dict(size=12,color=INK,_line=1.25)),
   ("衰退 −38.3% (p=0.0317)",dict(size=12.5,bold=True,color=RED,_line=1.25)),
   ("→ 原生獎勵已富梯度時，跨迭代記憶反而累積衝突 (Additive Shaping Conflict)。",
    dict(size=11,color=SLATE,_line=1.2))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,23)

# ============================================================
# SLIDE 24 — sparse unlock detail [orig S21]
# ============================================================
s=add_slide()
header(s,"解鎖極限：稀疏環境下的獎勵塑形",kicker="Sparse Win",kfill=GREEN,kw=1.7)
# CartPole
c1=box(s,0.7,1.65,5.85,2.05,fill=WHITE,line=HAIR,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
line_text(s,1.0,1.8,5.3,0.4,"CartPole-v1（存活挑戰）",size=15,color=INK,bold=True,latin=ZH,ea=ZH)
line_text(s,1.0,2.35,2.6,0.4,"原生 DQN (B0)",size=11,color=MUTE,latin=ZH,ea=ZH)
line_text(s,1.0,2.65,2.6,0.6,"154.8",size=26,color=MUTE,bold=True,latin=SERIF)
arrow(s,2.95,2.95,3.55,2.95,GREEN,2.5)
line_text(s,3.7,2.35,2.6,0.4,"Hermes-full (B3)",size=11,color=GREEN,latin=ZH,ea=ZH)
line_text(s,3.7,2.65,2.6,0.6,"334.4",size=26,color=GREEN,bold=True,latin=SERIF)
chip(s,4.95,3.35,"+116%",fill=GREEN,w=1.2,h=0.34,size=12)
line_text(s,1.0,3.32,3.5,0.3,"效能翻倍，解鎖原生無法解的任務",size=10.5,color=SLATE,latin=ZH,ea=ZH)
# MountainCar
c2=box(s,0.7,3.95,5.85,2.05,fill=WHITE,line=HAIR,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
line_text(s,1.0,4.1,5.3,0.4,"MountainCar-v0（動量挑戰）",size=15,color=INK,bold=True,latin=ZH,ea=ZH)
line_text(s,1.0,4.65,2.6,0.4,"原生 DQN (B0)",size=11,color=MUTE,latin=ZH,ea=ZH)
line_text(s,1.0,4.95,2.6,0.6,"−193.4",size=24,color=MUTE,bold=True,latin=SERIF)
arrow(s,3.05,5.25,3.65,5.25,GREEN,2.5)
line_text(s,3.8,4.65,2.6,0.4,"Hermes-full (B3)",size=11,color=GREEN,latin=ZH,ea=ZH)
line_text(s,3.8,4.95,2.6,0.6,"−132.5",size=24,color=GREEN,bold=True,latin=SERIF)
chip(s,5.0,5.6,"+31.5%",fill=GREEN,w=1.2,h=0.34,size=12)
line_text(s,1.0,5.62,3.8,0.3,"五個種子一致收斂（std=3.08）",size=10.5,color=SLATE,latin=ZH,ea=ZH)
# right takeaway
tb2=box(s,6.85,1.65,5.75,4.35,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
shape_text(tb2,[("為什麼稀疏環境會大勝？",dict(size=15,bold=True,color=GREEN,_sa=10)),
   ("原生獎勵幾乎不含學習訊號",dict(size=13,bold=True,color=INK,_line=1.2,_sa=2)),
   ("（僅二元存活 / 固定時間懲罰），原生 DQN 在預算內幾乎解不出題。",dict(size=12,color=SLATE,_line=1.25,_sa=10)),
   ("LLM 寫出的塑形 = 唯一明確子目標",dict(size=13,bold=True,color=INK,_line=1.2,_sa=2)),
   ("MountainCar 上 Gemma 寫出「每步 +0.5 動量」，媲美經典文獻的人類專家塑形。",dict(size=12,color=SLATE,_line=1.25,_sa=10)),
   ("B0 成功率：CartPole 0% · MountainCar 0% · Acrobot 62%",dict(size=11.5,bold=True,color=RED,_line=1.2))],
   anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT)
footer(s,24)

# ============================================================
# SLIDE 25 — variance signature [orig S22]
# ============================================================
s=add_slide()
header(s,"病理分析：變異性指紋 (Variance Signature)",kicker="Variance",kfill=SLATE,kw=1.7)
cd=CategoryChartData()
cd.categories=["MountainCar","Acrobot","LunarLander","CartPole"]
cd.add_series("per-seed std",(3.08,4.39,91.40,113.18))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(0.7),Inches(1.55),Inches(7.2),Inches(4.6),cd)
chart=gf.chart; style_chart(chart,legend=False,val_axis=True,cat_size=12)
data_labels(chart.plots[0],fmt='0.0',size=12)
ser=chart.series[0]
for i,c in enumerate((GREEN,GREEN,RED,AMBER)):
    ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb=c
line_text(s,0.7,6.2,7.2,0.3,"B3-hermes-full 各環境逐種子標準差（越低越穩定）",
          size=10,color=MUTE,align=PP_ALIGN.CENTER,latin=SANS,ea=ZH)
b1=box(s,8.15,1.65,4.5,1.7,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(b1,[("簡單物理 / 稀疏 → 極端穩定",dict(size=12.5,bold=True,color=GREEN,_sa=4)),
   ("MountainCar (3.08)、Acrobot (4.39) 緊湊收斂——單一明確目標，LLM 穩定收斂至近最佳唯一解。",
    dict(size=11.5,color=INK,_line=1.25))],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
b2=box(s,8.15,3.5,4.5,2.65,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(b2,[("豐富塑形空間 → 高變異風險",dict(size=12.5,bold=True,color=RED,_sa=4)),
   ("LunarLander (91.40) 高低落差極大。",dict(size=11.5,color=INK,_line=1.2,_sa=4)),
   ("致命的 seed_43（得分 11.6）：LLM 生成「重罰垂直速度」+「獎勵腳部接觸」的矛盾組合，",
    dict(size=11.5,color=SLATE,_line=1.25)),
   ("智能體學會「貼地懸停但不降落」的退化策略。",dict(size=11.5,bold=True,color=RED,_line=1.25))],
   anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,25)

# ============================================================
# SLIDE 26 — per-iteration trajectories (schematic) [orig S23]
# ============================================================
s=add_slide()
header(s,"軌跡透視：記憶機制如何引發干擾？",kicker="Trajectories",kfill=SLATE,kw=1.9)
# panel A: LunarLander chaos
pa=box(s,0.7,1.6,6.6,2.0,fill=WHITE,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
line_text(s,0.9,1.7,5.5,0.3,"LunarLander-v3（密集獎勵）",size=12,color=RED,bold=True,latin=ZH,ea=ZH)
ax,ay,aw,ah=1.4,3.4,5.6,0.0
connector(s,1.3,3.45,7.0,3.45,HAIR,1.0)
polyline(s,[(1.5,2.4),(2.4,2.15),(3.3,3.25),(4.2,2.0),(5.1,3.3),(6.0,2.25),(6.8,2.55)],RED,2.25)
line_text(s,5.3,2.95,1.8,0.3,"混沌震盪",size=12,color=RED,bold=True,latin=ZH,ea=ZH)
# panel B: MountainCar converge
pb=box(s,0.7,3.85,6.6,2.0,fill=WHITE,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
line_text(s,0.9,3.95,5.5,0.3,"MountainCar-v0（稀疏獎勵）",size=12,color=GREEN,bold=True,latin=ZH,ea=ZH)
connector(s,1.3,5.65,7.0,5.65,HAIR,1.0)
polyline(s,[(1.5,5.5),(2.4,5.25),(3.3,5.0),(4.2,4.78),(5.1,4.62),(6.0,4.52),(6.8,4.48)],GREEN,2.25)
line_text(s,5.2,4.5,1.9,0.3,"單調收斂",size=12,color=GREEN,bold=True,latin=ZH,ea=ZH)
line_text(s,0.9,5.55,5,0.3,"iteration 1 → 5",size=9,color=MUTE,latin=SANS,ea=ZH)
# right note
rb=box(s,7.6,1.6,5.0,4.25,fill=STICKY,line=STICKY_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
shape_text(rb,[("關鍵詮釋",dict(size=14,bold=True,color=AMBER,_sa=8)),
   ("當原生獎勵已有強烈梯度時，",dict(size=14,color=INK,bold=True,_line=1.3)),
   ("跨迭代記憶反而讓 LLM",dict(size=14,color=INK,bold=True,_line=1.3)),
   ("累積了「衝突的干擾項」。",dict(size=14,color=RED,bold=True,_line=1.3,_sa=10)),
   ("Additive Shaping Conflict",dict(size=12,italic=True,color=SLATE,_line=1.2,_sa=8)),
   ("（示意圖：依論文圖 5 之軌跡型態繪製）",dict(size=10,color=MUTE,_line=1.2))],
   anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT)
footer(s,26)

# ============================================================
# SLIDE 27 — Part 2 generalization (grouped bar) [orig S24]
# ============================================================
s=add_slide()
header(s,"Part 2 泛化驗證：這只是底層演算法的特例嗎？",kicker="Part 2",kfill=SLATE,kw=1.6)
line_text(s,0.7,1.4,11.9,0.35,"固定獎勵管線，僅切換 DQN 變體（vanilla / Double / Dueling）——Δ% = Hermes vs B0。",
          size=12.5,color=SLATE,latin=ZH,ea=ZH)
cd=CategoryChartData()
cd.categories=["LunarLander","CartPole","MountainCar","Acrobot"]
cd.add_series("vanilla",(-11.4,116.1,31.5,57.5))
cd.add_series("Double",(-23.4,113.4,31.0,65.3))
cd.add_series("Dueling",(-17.3,39.1,26.1,22.9))
gf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(0.7),Inches(1.85),Inches(8.2),Inches(4.4),cd)
chart=gf.chart; style_chart(chart,legend=True,val_axis=True,cat_size=11)
cols=[SLATE,BLUE,GREEN]
for i,c in enumerate(cols):
    chart.series[i].format.fill.solid(); chart.series[i].format.fill.fore_color.rgb=c
data_labels(chart.plots[0],fmt='0"%"',size=7.5,bold=False,color=SLATE)
rb=box(s,9.1,1.85,3.5,4.4,fill=GHOST,line=HAIR,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
shape_text(rb,[("三項確認",dict(size=14,bold=True,color=INK,_sa=10)),
   ("① 稀疏勝、密集反轉的型態",dict(size=12,bold=True,color=INK,_line=1.2)),
   ("在三種變體下完整複現。",dict(size=12,color=SLATE,_line=1.2,_sa=8)),
   ("② Hermes 在變體間無顯著差異",dict(size=12,bold=True,color=INK,_line=1.2)),
   ("(all p > 0.3)。",dict(size=12,color=SLATE,_line=1.2,_sa=8)),
   ("③ 證實模型無關性：獎勵設計",dict(size=12,bold=True,color=GREEN,_line=1.2)),
   ("的作用獨立於價值網路架構。",dict(size=12,color=SLATE,_line=1.2))],
   anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT)
footer(s,27)

# ============================================================
# SLIDE 28 — full results table (NEW backup)
# ============================================================
s=add_slide()
header(s,"完整結果數據表（附錄）",kicker="Appendix · Table 1",kfill=SLATE,kw=2.4)
line_text(s,0.7,1.4,11.9,0.35,"env_native_mean（n=5）｜粗綠 = 相對 B0 顯著勝出，紅斜 = 相對 no-memory 顯著敗退",
          size=11.5,color=MUTE,latin=ZH,ea=ZH)
heads=["條件","LunarLander (密集)","CartPole (稀疏)","MountainCar (稀疏)","Acrobot (稀疏)"]
cx=[0.7,3.7,6.05,8.4,10.75]; cw=[3.0,2.35,2.35,2.35,1.9]
y0=1.95
for h,x,w in zip(heads,cx,cw):
    line_text(s,x+0.05,y0,w-0.05,0.5,h,size=11.5,color=SLATE,bold=True,latin=ZH,ea=ZH,line=1.0)
# (value, style) style: 0 normal,1 green-bold,2 red-italic
rows=[("B0-env-native",[("173.22",0),("154.80",0),("−193.44",0),("−194.96",0)]),
      ("B1-handcrafted",[("77.77",0),("160.19",0),("−140.40",0),("−185.28",0)]),
      ("B2-gemma-oneshot",[("152.65",0),("187.64",0),("−153.09",0),("−83.21",0)]),
      ("B3-hermes-full",[("153.56",2),("334.44",1),("−132.53",1),("−82.92",0)]),
      ("B3-no-memory",[("248.77",1),("243.21",0),("−168.55",0),("−83.23",0)]),
      ("B3-no-AST",[("95.42",0),("220.81",0),("−134.59",0),("−83.58",0)])]
y=y0+0.65
for nm,vals in rows:
    hl = nm=="B3-hermes-full"
    if hl: box(s,0.65,y-0.04,11.95,0.6,fill=STICKY,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.1)
    else: connector(s,0.7,y+0.52,12.65,y+0.52,HAIR,0.75)
    line_text(s,cx[0]+0.05,y+0.06,cw[0],0.4,nm,size=12,color=INK,bold=hl,latin=MONO,ea=ZH)
    for (val,st),x,w in zip(vals,cx[1:],cw[1:]):
        col=INK; bold=False; ital=False
        if st==1: col=GREEN; bold=True
        elif st==2: col=RED; ital=True
        line_text(s,x+0.05,y+0.06,w-0.05,0.4,val,size=12.5,color=col,bold=bold,italic=ital,latin=SERIF,ea=ZH)
    y+=0.66
sticky(s,0.7,6.25,11.9,0.62,[("一句話：三個稀疏環境 B3-hermes-full 相對 B0 +31.5%~+116%；唯一的密集環境 LunarLander 反而被「記憶」拖累 −38.3%。",dict(size=12,bold=True,color=INK))],rotation=0)
footer(s,28)

# ============================================================
# SLIDE 29 — reward density hypothesis [orig S25]
# ============================================================
s=add_slide()
header(s,"結論法則：獎勵密度假說 (Reward Density Hypothesis)",kicker="Decision Rule",kfill=SLATE,kw=1.9)
q=box(s,0.7,2.7,3.1,1.6,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
shape_text(q,[("要不要使用",dict(size=14,bold=True,color=WHITE,_line=1.2)),
   ("記憶擴增的",dict(size=14,bold=True,color=WHITE,_line=1.2)),
   ("LLM 獎勵設計？",dict(size=14,bold=True,color=STICKY,_line=1.2))],anchor=MSO_ANCHOR.MIDDLE)
arrow(s,3.85,3.05,5.1,2.2,MUTE,2.0); arrow(s,3.85,3.95,5.1,4.85,MUTE,2.0)
# sparse branch
sb=box(s,5.15,1.5,3.4,1.55,fill=WHITE,line=GREEN_LN,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(sb,[("環境獎勵稀疏",dict(size=14,bold=True,color=GREEN,_sa=3)),
   ("單一目標 / 存活挑戰",dict(size=11.5,color=INK,_line=1.15)),
   ("原生梯度微弱、物理單純",dict(size=10.5,color=MUTE,_line=1.15))],anchor=MSO_ANCHOR.MIDDLE)
arrow(s,8.6,2.27,9.2,2.27,GREEN,2.5)
sr=box(s,9.25,1.5,3.4,1.55,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(sr,[("✓ 強力啟用",dict(size=15,bold=True,color=GREEN,_sa=3)),
   ("解鎖任務，效能 +30%~100%+，",dict(size=11,color=INK,_line=1.15)),
   ("且變異低。",dict(size=11,color=INK,_line=1.15))],anchor=MSO_ANCHOR.MIDDLE)
# dense branch
db=box(s,5.15,4.35,3.4,1.55,fill=WHITE,line=RED_LN,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(db,[("環境獎勵密集",dict(size=14,bold=True,color=RED,_sa=3)),
   ("多重連續塑形空間",dict(size=11.5,color=INK,_line=1.15)),
   ("原生梯度強烈、塑形豐富",dict(size=10.5,color=MUTE,_line=1.15))],anchor=MSO_ANCHOR.MIDDLE)
arrow(s,8.6,5.12,9.2,5.12,RED,2.5)
dr=box(s,9.25,4.35,3.4,1.55,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
shape_text(dr,[("⚠ 謹慎 / 禁用記憶",dict(size=14,bold=True,color=RED,_sa=3)),
   ("易產生加法干擾，引發高變異",dict(size=11,color=INK,_line=1.15)),
   ("與退化 (std 高達 91+)。",dict(size=11,color=INK,_line=1.15))],anchor=MSO_ANCHOR.MIDDLE)
footer(s,29)

# ============================================================
# SLIDE 30 — divider 5
# ============================================================
divider(5,"Conclusion","結論與未來展望")

# ============================================================
# SLIDE 31 — conclusion findings [orig S27]
# ============================================================
s=add_slide()
header(s,"Hermes-DQN：貢獻與核心發現",kicker="Conclusion",kfill=SLATE,kw=1.7)
line_text(s,0.7,1.5,11.9,0.7,"系統性整合「開源 LLM + 四層記憶 + AST 感知緩衝」，歷經 4 大環境、120 次嚴格訓練與消融實驗。",
          size=13.5,color=INK,latin=ZH,ea=ZH,line=1.25)
mods=[("開源 LLM","reward author",GREEN),("四層記憶","4-Tier Memory",BLUE),("AST 緩衝","AST-aware Buffer",AMBER)]
x=1.6
for t,en,c in mods:
    b=box(s,x,2.55,2.7,1.25,fill=WHITE,line=c,lw=1.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    shape_text(b,[(t,dict(size=16,bold=True,color=c,_sa=3)),(en,dict(size=10.5,color=MUTE))],anchor=MSO_ANCHOR.MIDDLE)
    if x<6: arrow(s,x+2.7,3.17,x+3.0,3.17,MUTE,2.5)
    x+=3.0
myth=box(s,0.7,4.35,11.9,1.85,fill=INK,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
shape_text(myth,[("核心發現：破除「記憶必定有益」的迷思",dict(size=17,bold=True,color=STICKY,_sa=8)),
   ("記憶對 LLM 獎勵設計的效益取決於「獎勵密度」——稀疏環境大幅受益，密集環境反而受害。",
    dict(size=13.5,color=WHITE,_line=1.3,_sa=3)),
   ("這是該領域首次提出、且具統計顯著的「任務相依性反轉」報告。",dict(size=13.5,bold=True,color=WHITE,_line=1.3))],
   anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.LEFT)
footer(s,31)

# ============================================================
# SLIDE 32 — sparse vs dense contrast [orig S28]
# ============================================================
s=add_slide()
header(s,"兩個世界：稀疏 vs 密集",kicker="Two Regimes",kfill=SLATE,kw=1.9)
# sparse column
sc=box(s,0.7,1.6,5.85,4.55,fill=GREEN_BG,line=GREEN_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
up=box(s,1.05,1.95,0.9,0.9,fill=GREEN,shape=MSO_SHAPE.UP_ARROW)
line_text(s,2.1,1.95,4.2,0.9,"稀疏獎勵\nSparse Environments",size=16,color=GREEN,bold=True,latin=SERIF,ea=ZH,anchor=MSO_ANCHOR.MIDDLE,line=1.1)
sp=[("任務","CartPole · MountainCar · Acrobot"),
    ("效能","+25% ~ +116% 顯著提升 (vs B0)"),
    ("特性","正向協同、變異極低 (std 3~4)"),
    ("結論","開源 LLM + 記憶發揮強大優勢")]
yy=3.2
for k,v in sp:
    line_text(s,1.1,yy,1.0,0.35,k,size=12,color=GREEN,bold=True,latin=ZH,ea=ZH)
    line_text(s,2.15,yy,4.25,0.6,v,size=12,color=INK,latin=ZH,ea=ZH,line=1.1)
    yy+=0.72
# dense column
dc=box(s,6.75,1.6,5.85,4.55,fill=RED_BG,line=RED_LN,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
dn=box(s,7.1,1.95,0.9,0.9,fill=RED,shape=MSO_SHAPE.DOWN_ARROW)
line_text(s,8.15,1.95,4.2,0.9,"密集獎勵\nDense Environments",size=16,color=RED,bold=True,latin=SERIF,ea=ZH,anchor=MSO_ANCHOR.MIDDLE,line=1.1)
dp=[("任務","LunarLander（唯一密集環境）"),
    ("效能","記憶導致效能衰退 −38%"),
    ("特性","高變異 (std 91+)、可能策略崩潰"),
    ("結論","該領域首次具統計顯著的負向報告")]
yy=3.2
for k,v in dp:
    line_text(s,7.15,yy,1.0,0.35,k,size=12,color=RED,bold=True,latin=ZH,ea=ZH)
    line_text(s,8.2,yy,4.25,0.6,v,size=12,color=INK,latin=ZH,ea=ZH,line=1.1)
    yy+=0.72
footer(s,32)

# ============================================================
# SLIDE 33 — limitations (expanded to 6, from paper)
# ============================================================
s=add_slide()
header(s,"研究限制 (Constraints & Limitations)",kicker="Limitations",kfill=SLATE,kw=1.8)
lims=[("小樣本規模","n=5 僅對大效應 (d≥1) 有力；部分記憶效應方向明確卻無法定論。後續以 n=10/20 重做。"),
      ("密集環境僅一個","「密度」與「塑形空間」在四環境面板中部分混淆；負向發現目前僅由 LunarLander 支撐，需 LunarLanderContinuous 分離兩變數。"),
      ("B1 為暫時基準","人工塑形由作者自撰，非第三方；不足以作為決定性人類基準。"),
      ("單一 LLM","變異可能與該模型抽樣特性有關；Llama 3.3 / Qwen 3 / DeepSeek-V3 重做為自然延伸。"),
      ("固定 5 輪迭代","閉環迴圈硬編碼為 5 輪；更長迭代能否逆轉負向效應未知。"),
      ("僅部分 Rainbow","只實作 Double + Dueling；PER / Multi-step / Noisy / 分佈式 Q 留待後續。")]
pos=[(0.7,1.65),(4.7,1.65),(8.7,1.65),(0.7,4.0),(4.7,4.0),(8.7,4.0)]
cw=3.85; ch=2.15
for (t,d),(x,y) in zip(lims,pos):
    b=box(s,x,y,cw,ch,fill=WHITE,line=HAIR,lw=1.25,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    box(s,x,y,cw,0.1,fill=AMBER,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.5)
    shape_text(b,[(t,dict(size=13.5,bold=True,color=INK,_sa=6)),(d,dict(size=11,color=SLATE,_line=1.25))],
               anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT)
footer(s,33)

# ============================================================
# SLIDE 34 — future roadmap [orig S30]
# ============================================================
s=add_slide()
header(s,"未來研究藍圖",kicker="Future Work",kfill=SLATE,kw=1.7)
steps=[("Step 1","連續控制擴張","進軍 MuJoCo 等高維連續動作基準。"),
       ("Step 2","勢能約束生成","將理論邊界注入 LLM Prompt，從數學上保證生成的塑形獎勵不破壞原始最優策略。"),
       ("Step 3","自適應迭代機制","捨棄固定輪數，改以 Fitness 曲線動態判定是否終止迭代迴圈。"),
       ("Step 4","多模型基準整合","引入多種開源模型橫向 Benchmark，隔離出模型特定的偏誤特徵。")]
# staircase
n=len(steps); baseY=5.9; stepH=0.95; x0=0.9; xw=2.95
for i,(st,t,d) in enumerate(steps):
    y=baseY-i*stepH
    dot(s,x0+i*xw+0.15,y,0.12,AMBER)
    if i<n-1:
        connector(s,x0+i*xw+0.15,y,x0+(i+1)*xw+0.15,y-stepH,SLATE,2.0)
    b=box(s,x0+i*xw-0.05,y-stepH-0.7,xw-0.1,1.35,fill=WHITE,line=HAIR,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    shape_text(b,[(st,dict(size=11,bold=True,color=AMBER,_sa=2)),
                  (t,dict(size=13.5,bold=True,color=INK,_sa=4)),
                  (d,dict(size=10.5,color=SLATE,_line=1.2))],anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT)
line_text(s,0.7,6.55,11.9,0.4,"陳盛茂 ・ 林仙安 ・ 辛語柔 ・ 陳冠宇　｜　國立中興大學 資訊管理學研究所",
          size=12,color=MUTE,align=PP_ALIGN.CENTER,latin=ZH,ea=ZH)
footer(s,34)

# ---------------- save ----------------
out=r"C:\Users\Mao\Desktop\DRL\Final Project\PPT\PPT_第三版.pptx"
prs.save(out)
print("SAVED:",out,"slides=",len(prs.slides._sldIdLst))
